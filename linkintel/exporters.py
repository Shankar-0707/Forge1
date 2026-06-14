import os
from xhtml2pdf import pisa
from pptx import Presentation
from pptx.util import Inches

def export_pdf(report_obj: dict, output_path: str):
    """Generate a PDF report using xhtml2pdf."""
    s = report_obj.get("summary", {})
    
    html = f"""
    <html>
    <head>
    <style>
        body {{ font-family: Helvetica, sans-serif; font-size: 12px; color: #333; }}
        h1 {{ color: #2c3e50; border-bottom: 2px solid #2c3e50; padding-bottom: 5px; }}
        h2 {{ color: #34495e; margin-top: 20px; }}
        table {{ width: 100%; border-collapse: collapse; margin-bottom: 20px; }}
        th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
        th {{ background-color: #f2f2f2; font-weight: bold; }}
        .metric {{ font-weight: bold; color: #e74c3c; }}
    </style>
    </head>
    <body>
        <h1>Internal Linking Intelligence Report</h1>
        <p><strong>Site:</strong> {report_obj.get('site', 'Unknown')}</p>
        <p><strong>Pages Crawled:</strong> {s.get('pages_crawled', 0)}</p>
        
        <h2>1. Executive Summary</h2>
        <table>
            <tr><th>Metric</th><th>Value</th></tr>
            <tr><td>Indexable Pages</td><td>{s.get('indexable_pages', 0)}</td></tr>
            <tr><td>Internal Links</td><td>{s.get('internal_links', 0)}</td></tr>
            <tr><td>Orphan Pages</td><td><span class="metric">{s.get('orphan_pages', 0)}</span></td></tr>
            <tr><td>Broken Internal Links</td><td><span class="metric">{s.get('broken_internal_links', 0)}</span></td></tr>
            <tr><td>Generic Anchors</td><td>{s.get('generic_anchors', 0)}</td></tr>
            <tr><td>Topical Clusters</td><td>{s.get('topical_clusters', 0)}</td></tr>
            <tr><td>Link Recommendations</td><td>{s.get('link_recommendations', 0)}</td></tr>
        </table>
        
        <h2>2. Contextual Link Recommendations</h2>
        <table>
            <tr><th>Source Page</th><th>Target Page</th><th>Suggested Anchor</th></tr>
    """
    recs = report_obj.get("link_recommendations", [])
    for rec in recs[:20]:
        html += f"<tr><td>{rec.get('source', '')[-40:]}</td><td>{rec.get('target', '')[-40:]}</td><td>{rec.get('suggested_anchor', '')}</td></tr>"
        
    html += """
        </table>
        
        <h2>3. Topical Clusters</h2>
        <table>
            <tr><th>Cluster Name</th><th>Size</th><th>Authority</th><th>Keywords</th></tr>
    """
    clusters = report_obj.get("topical_clusters", [])
    for cl in clusters[:15]:
        html += f"<tr><td>{cl.get('name', cl.get('key'))}</td><td>{cl.get('size', 0)}</td><td>{cl.get('authority', '')}</td><td>{', '.join(cl.get('keywords', [])[:3])}</td></tr>"
        
    html += """
        </table>
    </body>
    </html>
    """
    
    with open(output_path, "wb") as f:
        pisa.CreatePDF(html, dest=f)


def export_pptx(report_obj: dict, output_path: str):
    """Generate a PPTX report using python-pptx."""
    prs = Presentation()
    
    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Internal Linking Intelligence"
    subtitle.text = f"Site: {report_obj.get('site', 'Unknown')}\nPages Crawled: {report_obj.get('pages_crawled', 0)}"
    
    s = report_obj.get("summary", {})
    
    # Slide 2: Executive Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes2 = slide2.shapes
    shapes2.title.text = "Executive Summary"
    tf2 = shapes2.placeholders[1].text_frame
    tf2.text = f"Indexable Pages: {s.get('indexable_pages', 0)}"
    tf2.add_paragraph().text = f"Internal Links: {s.get('internal_links', 0)}"
    tf2.add_paragraph().text = f"Orphan Pages: {s.get('orphan_pages', 0)}"
    tf2.add_paragraph().text = f"Broken Internal Links: {s.get('broken_internal_links', 0)}"
    tf2.add_paragraph().text = f"Topical Clusters: {s.get('topical_clusters', 0)}"
    tf2.add_paragraph().text = f"Link Recommendations: {s.get('link_recommendations', 0)}"
    
    # Slide 3: Top Link Recommendations
    slide3 = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    slide3.shapes.title.text = "Top Link Recommendations"
    
    recs = report_obj.get("link_recommendations", [])
    rows = min(len(recs) + 1, 6) # Max 5 rows + header
    cols = 3
    if rows > 1:
        table = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table
        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(3.0)
        table.columns[2].width = Inches(3.0)
        
        table.cell(0, 0).text = "Source"
        table.cell(0, 1).text = "Target"
        table.cell(0, 2).text = "Suggested Anchor"
        
        for i in range(1, rows):
            rec = recs[i-1]
            table.cell(i, 0).text = rec.get("source", "")[-30:]
            table.cell(i, 1).text = rec.get("target", "")[-30:]
            table.cell(i, 2).text = rec.get("suggested_anchor", "")
            
    prs.save(output_path)
